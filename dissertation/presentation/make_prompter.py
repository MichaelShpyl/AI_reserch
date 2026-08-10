"""Build a rehearsal teleprompter from the deck.

Reading the talk track off a printed page while clicking through slides does not rehearse the
thing you actually have to do, which is stay on pace. This makes a single self-contained page that
shows one script at a time with a per-slide countdown and a running clock against the twenty-minute
limit, so you can find out where you drift before Vini's camera does.

    python dissertation/presentation/make_prompter.py

Writes dissertation/presentation/Rehearsal_Prompter.html. Open it in a browser, press space.
Everything is read out of the deck's speaker notes, so it cannot go stale against the slides.
"""
from __future__ import annotations

import json
import re
from pathlib import Path

from pptx import Presentation

ROOT = Path(__file__).resolve().parents[2]
PRES = ROOT / "dissertation" / "presentation"
DECK = PRES / "PreFinal_Presentation_Shpyl.pptx"
OUT = PRES / "Rehearsal_Prompter.html"
LIMIT = 20 * 60


def slides():
    prs = Presentation(DECK)
    rows = []
    for i, s in enumerate(prs.slides, 1):
        note = s.notes_slide.notes_text_frame.text
        m = re.match(r"\[([A-Z]+)[^\]]*\|\s*(\d+)s\]\s*(.*)", note, re.S)
        if not m:
            continue
        tier = {"CORE": "core", "DETAIL": "detail", "APPENDIX": "appendix"}[m.group(1)]
        eyebrow = title = ""
        for sh in s.shapes:
            if not sh.has_text_frame:
                continue
            t = sh.text_frame.text.strip()
            if not t or t in ("detail", "appendix"):
                continue
            if not eyebrow and t.isupper():
                eyebrow = t
                continue
            if not title and len(t) > 4:
                title = t.split("\n")[0]
        rows.append({"n": i, "tier": tier, "secs": int(m.group(2)),
                     "eyebrow": eyebrow, "title": title, "talk": m.group(3).strip()})
    return rows


HTML = """<!doctype html>
<meta charset="utf-8">
<title>Rehearsal prompter</title>
<style>
:root{--paper:#fbfaf8;--ink:#222831;--teal:#2b6777;--rust:#a63d2e;--muted:#5b6770;--line:#dfe3e6}
@media (prefers-color-scheme: dark){
  :root{--paper:#161a1d;--ink:#e8eaec;--teal:#7fb3c0;--rust:#e08b78;--muted:#9aa5ac;--line:#2b3236}
}
*{box-sizing:border-box}
body{margin:0;background:var(--paper);color:var(--ink);
  font:16px/1.6 Calibri,system-ui,sans-serif;height:100vh;display:flex;flex-direction:column}
header{display:flex;align-items:baseline;gap:1.2rem;padding:.7rem 1.6rem;
  border-bottom:1px solid var(--line);flex:none}
.slideno{font-weight:700;font-size:1.1rem}
.tier{font-size:.7rem;font-weight:700;letter-spacing:.08em;text-transform:uppercase;
  padding:.14rem .55rem;border-radius:.8rem;background:var(--line);color:var(--teal)}
.tier.detail{color:var(--muted)} .tier.appendix{color:var(--rust)}
.eyebrow{font-size:.72rem;letter-spacing:.08em;text-transform:uppercase;color:var(--teal);font-weight:700}
.clocks{margin-left:auto;display:flex;gap:1.6rem;align-items:baseline;
  font-family:Consolas,monospace;font-variant-numeric:tabular-nums}
.clocks b{font-size:1.5rem;font-weight:700}
.clocks .lbl{font-size:.66rem;letter-spacing:.07em;text-transform:uppercase;color:var(--muted);
  font-family:Calibri,sans-serif;font-weight:700;display:block}
#slideclock.over{color:var(--rust)}
#total.risk{color:var(--rust)}
main{flex:1;overflow-y:auto;padding:1.4rem 1.6rem 3rem;max-width:56rem}
h1{font-size:1.75rem;line-height:1.25;margin:.2rem 0 1.1rem;font-weight:700;text-wrap:balance}
.talk{font-size:1.5rem;line-height:1.62}
.talk .cue{color:var(--rust);font-weight:700}
.bar{height:5px;background:var(--line);flex:none}
.bar i{display:block;height:100%;width:0;background:var(--teal);transition:width .25s linear}
.bar i.over{background:var(--rust)}
footer{flex:none;padding:.55rem 1.6rem;border-top:1px solid var(--line);font-size:.8rem;
  color:var(--muted);display:flex;gap:1.4rem;flex-wrap:wrap}
kbd{font-family:Consolas,monospace;background:var(--line);border-radius:.25rem;padding:.05rem .35rem}
.paused{opacity:.45}
@media (max-width:860px){ .talk{font-size:1.2rem} h1{font-size:1.3rem} }
</style>
<header>
  <span class="slideno" id="num"></span>
  <span class="tier" id="tier"></span>
  <span class="eyebrow" id="eyebrow"></span>
  <span class="clocks">
    <span><span class="lbl">this slide</span><b id="slideclock">0:00</b></span>
    <span><span class="lbl">target</span><b id="target">0:00</b></span>
    <span><span class="lbl">elapsed</span><b id="total">0:00</b></span>
    <span><span class="lbl">to 20:00</span><b id="left">20:00</b></span>
  </span>
</header>
<div class="bar"><i id="fill"></i></div>
<main><h1 id="title"></h1><div class="talk" id="talk"></div></main>
<footer>
  <span><kbd>space</kbd> next slide</span>
  <span><kbd>&larr;</kbd> back</span>
  <span><kbd>p</kbd> pause</span>
  <span><kbd>r</kbd> restart</span>
  <span><kbd>d</kbd> include detail slides</span>
  <span id="mode"></span>
</footer>
<script>
const ALL = __DATA__;
const LIMIT = __LIMIT__;
let includeDetail = false, i = 0, running = false, paused = false;
let slideStart = 0, runStart = 0, pausedAt = 0, offset = 0;

const deck = () => includeDetail ? ALL : ALL.filter(s => s.tier === "core");
const $ = s => document.querySelector(s);
const mmss = t => `${Math.floor(t / 60)}:${String(Math.floor(t % 60)).padStart(2, "0")}`;

function render() {
  const d = deck(), s = d[Math.min(i, d.length - 1)];
  $("#num").textContent = `${i + 1} / ${d.length}`;
  $("#tier").textContent = s.tier;
  $("#tier").className = "tier " + s.tier;
  $("#eyebrow").textContent = s.eyebrow || "";
  $("#title").textContent = s.title || "";
  // Stage cues sit in square brackets in the script. They are instructions, not words to say.
  $("#talk").innerHTML = s.talk
    .replace(/[&<>]/g, c => ({ "&": "&amp;", "<": "&lt;", ">": "&gt;" }[c]))
    .replace(/\\[([^\\]]+)\\]/g, '<span class="cue">[$1]</span>');
  $("#target").textContent = mmss(s.secs);
  $("#mode").textContent = includeDetail
    ? `all tiers, ${mmss(d.reduce((a, x) => a + x.secs, 0))} of material`
    : `core path, target ${mmss(d.reduce((a, x) => a + x.secs, 0))}`;
  window.scrollTo(0, 0);
  document.querySelector("main").scrollTop = 0;
}

function tick() {
  if (!running || paused) return;
  const now = performance.now() / 1000;
  const d = deck(), s = d[Math.min(i, d.length - 1)];
  const onSlide = now - slideStart;
  const elapsed = now - runStart - offset;
  $("#slideclock").textContent = mmss(onSlide);
  $("#slideclock").classList.toggle("over", onSlide > s.secs);
  $("#total").textContent = mmss(elapsed);
  $("#left").textContent = (LIMIT - elapsed >= 0 ? "" : "-") + mmss(Math.abs(LIMIT - elapsed));
  $("#total").classList.toggle("risk", elapsed > LIMIT);
  const f = $("#fill");
  f.style.width = Math.min(100, onSlide / s.secs * 100) + "%";
  f.classList.toggle("over", onSlide > s.secs);
}
setInterval(tick, 100);

function advance(step) {
  const d = deck();
  const now = performance.now() / 1000;
  if (!running) { running = true; runStart = now; offset = 0; }
  i = Math.max(0, Math.min(d.length - 1, i + step));
  slideStart = now;
  render();
}

addEventListener("keydown", e => {
  if (e.key === " " || e.key === "ArrowRight") { e.preventDefault(); advance(1); }
  else if (e.key === "ArrowLeft") { e.preventDefault(); advance(-1); }
  else if (e.key === "p") {
    paused = !paused;
    document.body.classList.toggle("paused", paused);
    const now = performance.now() / 1000;
    if (paused) pausedAt = now; else { offset += now - pausedAt; slideStart += now - pausedAt; }
  } else if (e.key === "r") {
    running = false; paused = false; i = 0; offset = 0;
    document.body.classList.remove("paused");
    $("#slideclock").textContent = $("#total").textContent = "0:00";
    $("#left").textContent = mmss(LIMIT); $("#fill").style.width = "0";
    render();
  } else if (e.key === "d") {
    includeDetail = !includeDetail; i = 0; render();
  }
});
render();
</script>
"""


def main():
    rows = slides()
    html = (HTML.replace("__DATA__", json.dumps(rows, ensure_ascii=False))
                .replace("__LIMIT__", str(LIMIT)))
    OUT.write_text(html, encoding="utf-8")
    core = sum(r["secs"] for r in rows if r["tier"] == "core")
    print(f"wrote {OUT}")
    print(f"{len(rows)} slides, core path {core // 60}m{core % 60:02d}s")


if __name__ == "__main__":
    main()
