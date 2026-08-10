"""Rehearse the presentation and measure what it actually takes to say.

Every timing in the talk track so far has been words divided by an assumed speaking rate. That is a
decent estimate and it is still an estimate: it treats a slide of short declarative sentences the
same as a slide of one long clause, when the first takes noticeably longer to deliver because of
where you breathe.

This speaks every script with the Windows speech engine, measures the audio, and reports the real
duration against the allotted one. The engine is not a person, so the absolute rate is calibrated:
the whole core path is scaled to a chosen target pace and the per-slide distribution is what the
measurement contributes. Punctuation, sentence length and paragraph breaks all move that
distribution, and none of them show up in a word count.

    python dissertation/presentation/rehearse.py                 measure and report
    python dissertation/presentation/rehearse.py --audio core    also write a rehearsal track

The audio is a single file of the core path with a short gap at each slide change, for practising
against. It is deliberately not committed.
"""
from __future__ import annotations

import argparse
import contextlib
import re
import shutil
import subprocess
import sys
import tempfile
import wave
from pathlib import Path

from pptx import Presentation

ROOT = Path(__file__).resolve().parents[2]
PRES = ROOT / "dissertation" / "presentation"
DECK = PRES / "PreFinal_Presentation_Shpyl.pptx"
VOICE = "Microsoft Hazel Desktop"    # en-GB, closer to the room than the US voice
RATE = -2                            # slower than default; the calibration below does the rest
TARGET_WPM = 135                     # a careful presenting pace, with room to breathe
GAP_SECONDS = 1.2                    # silence between slides in the rehearsal track


CUE = re.compile(r"\s*\[[^\]]*\]")


def spoken(talk: str) -> str:
    """The words that are actually said. Anything in square brackets is a stage cue for the
    presenter, so it is neither spoken nor counted, and the calibration below depends on those two
    staying consistent with each other."""
    return CUE.sub("", talk).strip()


def scripts():
    """(number, tier, allotted seconds, title, script) for every slide, in order."""
    prs = Presentation(DECK)
    out = []
    for i, s in enumerate(prs.slides, 1):
        note = s.notes_slide.notes_text_frame.text
        m = re.match(r"\[([A-Z]+)[^\]]*\|\s*(\d+)s\]\s*(.*)", note, re.S)
        if not m:
            continue
        tier = {"CORE": "core", "DETAIL": "detail", "APPENDIX": "appendix"}[m.group(1)]
        title = ""
        for sh in s.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                t = sh.text_frame.text.strip().split("\n")[0]
                if t.isupper() or len(t) < 5 or t in ("detail", "appendix"):
                    continue
                title = t
                break
        out.append((i, tier, int(m.group(2)), title, m.group(3).strip()))
    return out


def speak_all(items, workdir: Path):
    """Synthesise every script to its own wav. One PowerShell process for the whole batch, because
    starting the speech engine costs more than the speaking does on the short ones."""
    manifest = workdir / "jobs.txt"
    lines = []
    for n, _, _, _, talk in items:
        # The script is written for a person, so the engine gets it verbatim. Tabs separate the
        # fields; the scripts contain none.
        lines.append(f"{n}\t{talk}")
    manifest.write_text("\n".join(lines), encoding="utf-8")

    ps = f"""
$ErrorActionPreference = 'Stop'
Add-Type -AssemblyName System.Speech
$synth = New-Object System.Speech.Synthesis.SpeechSynthesizer
try {{ $synth.SelectVoice('{VOICE}') }} catch {{ }}
$synth.Rate = {RATE}
# Prime the engine into a throwaway file. Without this the first real script carries the engine's
# start-up silence, which showed up as the title slide measuring seven seconds longer than it does.
$synth.SetOutputToWaveFile((Join-Path '{workdir.as_posix()}' 'prime.wav'))
$synth.Speak('Priming the speech engine.')
foreach ($line in [System.IO.File]::ReadAllLines('{manifest.as_posix()}', [System.Text.Encoding]::UTF8)) {{
  if (-not $line) {{ continue }}
  $i = $line.IndexOf([char]9)
  $n = $line.Substring(0, $i)
  $t = $line.Substring($i + 1)
  $out = Join-Path '{workdir.as_posix()}' ("s{{0:d3}}.wav" -f [int]$n)
  $synth.SetOutputToWaveFile($out)
  $synth.Speak($t)
}}
$synth.SetOutputToNull()
$synth.Dispose()
Write-Output 'done'
"""
    script = workdir / "speak.ps1"
    script.write_text(ps, encoding="utf-8")
    r = subprocess.run(["powershell", "-NoProfile", "-ExecutionPolicy", "Bypass",
                        "-File", str(script)], capture_output=True, text=True)
    if "done" not in r.stdout:
        sys.exit(f"speech synthesis failed:\n{r.stderr[:2000]}")


def duration(path: Path) -> float:
    with contextlib.closing(wave.open(str(path), "rb")) as w:
        return w.getnframes() / float(w.getframerate())


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--audio", choices=["core", "all"], help="also write a rehearsal track")
    args = ap.parse_args()

    items = scripts()
    work = Path(tempfile.mkdtemp(prefix="rehearse_"))
    try:
        print(f"speaking {len(items)} scripts, this takes a few minutes")
        speak_all(items, work)

        raw = {}
        for n, _, _, _, _ in items:
            f = work / f"s{n:03d}.wav"
            if not f.exists():
                sys.exit(f"no audio for slide {n}")
            raw[n] = duration(f)

        # Calibrate: scale so the core path lands at the target pace. Everything after this is the
        # engine's per-slide distribution rather than its absolute speed.
        core = [it for it in items if it[1] == "core"]
        core_words = sum(len(spoken(it[4]).split()) for it in core)
        core_raw = sum(raw[it[0]] for it in core)
        target_core = core_words / TARGET_WPM * 60
        k = target_core / core_raw
        print(f"engine pace {core_words / (core_raw / 60):.0f} wpm, scaled by {k:.3f} "
              f"to {TARGET_WPM} wpm\n")

        rows = []
        for n, tier, allotted, title, talk in items:
            measured = raw[n] * k
            rows.append((n, tier, allotted, measured, title, len(spoken(talk).split())))

        # The engine is not a person. Proper nouns it cannot pronounce, and long runs of digits,
        # make it crawl in a way a presenter would not, so a slide whose measured length is far from
        # what its word count implies is flagged rather than trusted.
        def by_words(w):
            return w / TARGET_WPM * 60

        print(f"{'#':>3} {'tier':9} {'set':>5} {'said':>6} {'diff':>6}  slide")
        outliers = []
        for n, tier, allotted, measured, title, words in rows:
            d = measured - allotted
            ratio = measured / by_words(words) if words else 1.0
            if not (0.78 <= ratio <= 1.30):
                outliers.append((n, title, ratio, by_words(words)))
                mark = "  ?? engine outlier"
            else:
                mark = "  <<" if abs(d) >= 8 else ""
            print(f"{n:3d} {tier:9} {allotted:4d}s {measured:5.1f}s {d:+5.1f}s{mark}  {title[:44]}")

        print()
        for tier in ("core", "detail", "appendix"):
            g = [r for r in rows if r[1] == tier]
            a = sum(r[2] for r in g)
            m = sum(r[3] for r in g)
            print(f"{tier:9} set {a // 60}m{a % 60:02d}s   said {int(m) // 60}m{int(m) % 60:02d}s"
                  f"   {m - a:+.0f}s")
        core_m = sum(r[3] for r in rows if r[1] == "core")
        print(f"\nCORE against the 20 minute limit: {int(core_m) // 60}m{int(core_m) % 60:02d}s, "
              f"{1200 - core_m:.0f}s of slack")

        if outliers:
            print(f"\n{len(outliers)} slides where the engine disagrees with the word count. "
                  "Trust the word count on these:")
            for n, title, ratio, est in outliers:
                print(f"  {n:3d} engine {ratio:.2f}x the word-count estimate of {est:.0f}s   "
                      f"{title[:44]}")

        flagged = {o[0] for o in outliers}
        big = [r for r in rows if abs(r[3] - r[2]) >= 8 and r[0] not in flagged]
        if big:
            print(f"\n{len(big)} slides more than 8 seconds out. Suggested seconds:")
            for n, tier, allotted, measured, title, _ in big:
                print(f"  {n:3d} {tier:9} {allotted:4d} -> {round(measured):4d}   {title[:44]}")

        if args.audio:
            want = [it for it in items if args.audio == "all" or it[1] == "core"]
            dest = PRES / f"Rehearsal_{args.audio}.mp3"
            # Speed the track up by the same factor the report is calibrated with, so what you
            # hear is the pace you are aiming for rather than the engine's default crawl.
            join(work, [it[0] for it in want], dest, speed=1.0 / k)
            print(f"\nwrote {dest}")
    finally:
        shutil.rmtree(work, ignore_errors=True)


def join(work: Path, numbers, dest: Path, speed: float = 1.0):
    """Concatenate the per-slide wavs with a gap between them, and encode to mp3 at the target
    pace. The gaps are inserted before the tempo change, so they scale with everything else."""
    import imageio_ffmpeg
    listing = work / "list.txt"
    silence = work / "gap.wav"
    with contextlib.closing(wave.open(str(work / f"s{numbers[0]:03d}.wav"), "rb")) as w:
        params = w.getparams()
    with contextlib.closing(wave.open(str(silence), "wb")) as w:
        w.setparams(params)
        w.writeframes(b"\x00" * int(params.framerate * params.sampwidth *
                                    params.nchannels * GAP_SECONDS))
    lines = []
    for n in numbers:
        lines.append(f"file '{(work / f's{n:03d}.wav').as_posix()}'")
        lines.append(f"file '{silence.as_posix()}'")
    listing.write_text("\n".join(lines), encoding="utf-8")
    cmd = [imageio_ffmpeg.get_ffmpeg_exe(), "-y", "-f", "concat", "-safe", "0", "-i", str(listing)]
    if abs(speed - 1.0) > 0.01:
        cmd += ["-filter:a", f"atempo={speed:.4f}"]
    cmd += ["-c:a", "libmp3lame", "-b:a", "96k", str(dest)]
    subprocess.run(cmd, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, check=True)


if __name__ == "__main__":
    main()
