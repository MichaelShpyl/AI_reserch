"""Measure the submission video's talk track against the 10-minute ceiling.

Same method as rehearse.py, pointed at the submission deck. Every script is spoken by the Windows
speech engine, the audio measured, and the whole thing calibrated to a target presenting pace, so
the total is what the words take to say rather than a word count divided by an assumption.

The ceiling is not advisory. The department's guidance says examiners stop reviewing at 10:00, so a
script that runs to 10:30 loses its last thirty seconds rather than overrunning politely.

    python dissertation/presentation/rehearse_submission.py
"""

from __future__ import annotations

import contextlib
import re
import subprocess
import sys
import tempfile
import wave
from pathlib import Path

from pptx import Presentation

HERE = Path(__file__).resolve().parent
DECK = HERE / "Submission_Video_Deck_Shpyl.pptx"
LIMIT = 600                          # 10:00, hard
VOICE = "Microsoft Zira Desktop"
RATE = -2
TARGET_WPM = 135
CUE = re.compile(r"\[.*?\]", re.S)


def spoken(talk: str) -> str:
    return CUE.sub("", talk).strip()


def scripts():
    """(number, allotted seconds, title, script) for every slide."""
    prs = Presentation(str(DECK))
    out = []
    for i, s in enumerate(prs.slides, 1):
        note = s.notes_slide.notes_text_frame.text
        m = re.match(r"(\d+)\.\s*(.*?)\s*\((\d+)s\)\s*(.*)", note, re.S)
        if not m:
            continue
        out.append((i, int(m.group(3)), m.group(2).strip(), m.group(4).strip()))
    return out


def speak_all(items, workdir: Path):
    manifest = workdir / "jobs.txt"
    manifest.write_text(
        "\n".join(f"{n}\t{spoken(talk)}" for n, _, _, talk in items), encoding="utf-8")
    ps = f"""
$ErrorActionPreference = 'Stop'
Add-Type -AssemblyName System.Speech
$synth = New-Object System.Speech.Synthesis.SpeechSynthesizer
try {{ $synth.SelectVoice('{VOICE}') }} catch {{ }}
$synth.Rate = {RATE}
$synth.SetOutputToWaveFile((Join-Path '{workdir.as_posix()}' 'prime.wav'))
$synth.Speak('Priming the speech engine.')
foreach ($line in [System.IO.File]::ReadAllLines('{manifest.as_posix()}', [System.Text.Encoding]::UTF8)) {{
  if (-not $line) {{ continue }}
  $i = $line.IndexOf([char]9)
  $n = $line.Substring(0, $i)
  $t = $line.Substring($i + 1)
  $out = Join-Path '{workdir.as_posix()}' ("s{{0:d3}}.wav" -f [int]$n)
  $synth.SetOutputToWaveFile($out)
  $synth.Speak($t)
}}
$synth.SetOutputToNull()
$synth.Dispose()
Write-Output 'done'
"""
    script = workdir / "speak.ps1"
    script.write_text(ps, encoding="utf-8")
    r = subprocess.run(["powershell", "-NoProfile", "-ExecutionPolicy", "Bypass",
                        "-File", str(script)], capture_output=True, text=True)
    if "done" not in r.stdout:
        sys.exit(f"speech synthesis failed:\n{r.stderr[:2000]}")


def duration(path: Path) -> float:
    with contextlib.closing(wave.open(str(path), "rb")) as w:
        return w.getnframes() / float(w.getframerate())


def main() -> None:
    if not DECK.exists():
        sys.exit(f"Missing {DECK}")
    items = scripts()
    if not items:
        sys.exit("No parsable speaker notes found in the deck.")

    words = sum(len(spoken(t).split()) for _, _, _, t in items)
    target_total = words / TARGET_WPM * 60

    with tempfile.TemporaryDirectory() as td:
        work = Path(td)
        speak_all(items, work)
        raw = {n: duration(work / f"s{n:03d}.wav") for n, _, _, _ in items}

    engine_total = sum(raw.values())
    scale = target_total / engine_total            # calibrate the engine to a human pace
    print(f"{len(items)} slides, {words:,} spoken words")
    print(f"calibrated to {TARGET_WPM} wpm (engine scaled by {scale:.3f})\n")
    print(f"{'#':>3} {'set':>5} {'said':>7} {'delta':>7}  slide")

    total = 0.0
    flags = []
    for n, secs, title, _ in items:
        said = raw[n] * scale
        total += said
        d = said - secs
        mark = "  <<" if abs(d) > 8 else ""
        print(f"{n:>3} {secs:>4}s {said:>6.1f}s {d:>+6.1f}s{mark}  {title[:44]}")
        if abs(d) > 8:
            flags.append((n, secs, round(said), title))

    m, s = int(total) // 60, int(total) % 60
    print(f"\nTOTAL  set {sum(i[1] for i in items)}s   said {m}:{s:02d}")
    if total > LIMIT:
        over = total - LIMIT
        print(f"\n  OVER THE 10 MINUTE LIMIT BY {int(over)}s. "
              f"Examiners stop at 10:00, so that is {int(over)}s nobody watches.")
        print(f"  Cut about {int(over / 60 * TARGET_WPM) + 15} words.")
    else:
        print(f"\n  Inside the limit with {int(LIMIT - total)}s to spare.")

    if flags:
        print("\nSlides more than 8s from their stated time. Suggested seconds:")
        for n, secs, said, title in flags:
            print(f"    {n:>3} {secs:>4} -> {said:>4}   {title[:48]}")


if __name__ == "__main__":
    main()
