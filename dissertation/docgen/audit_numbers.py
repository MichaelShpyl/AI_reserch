"""Check every number in the write-up and the deck against the results files.

The existing consistency audit checks a hand-maintained list of headline figures. That list only
contains numbers I remembered to add to it, which is exactly the wrong property for a safety net.
The summary claim that the hybrid cuts false accusations "by a factor of three to eight" sat in
three places for weeks and was wrong: dividing the per-domain rates in hybrid_fusion.json gives
3.3, 3.4, 3.9 and 4.9. Nothing caught it because nobody had thought to check it.

So this works the other way round. It pulls every number out of the prose and the slides, pulls
every number out of outputs/*.json, and reports the ones in the text that no results file supports.

It cannot know what a number means, so it will flag things that are fine: sample sizes stated in
prose, years, section numbers, figures quoted from other people's papers. The output is a list to
read, not a list of defects. What matters is that a number invented or mistyped anywhere in 34,000
words cannot hide from it.

    python dissertation/docgen/audit_numbers.py            the chapters
    python dissertation/docgen/audit_numbers.py --deck     the chapters and the presentation
"""
from __future__ import annotations

import argparse
import json
import re
from pathlib import Path

ROOT = Path(__file__).resolve().parents[2]
CHAPTERS = ROOT / "dissertation" / "chapters"
OUTPUTS = ROOT / "outputs"
DECK = ROOT / "dissertation" / "presentation" / "PreFinal_Presentation_Shpyl.pptx"

# Numbers that are not claims about results and would otherwise dominate the report.
IGNORE_EXACT = {
    "0", "1", "2", "3", "4", "5", "6", "7", "8", "9", "10", "11", "12",
    "100", "1000", "2020", "2021", "2022", "2023", "2024", "2025", "2026",
    "0.0", "1.0", "0.5",
}
# Contexts where a number is structural rather than empirical.
SKIP_LINE = re.compile(
    r"^\s*(#|\||!\[|```)"                    # headings, tables handled separately, figures, code
    r"|^\s*(Chapter|Section|Figure|Table|Code Listing)\s+\d",
    re.I)
NEAR_STRUCTURAL = re.compile(
    r"(Section|Chapter|Figure|Table|Listing|Appendix|Objective|page|Meeting|"
    r"RTX|GB|GPU|version|v\d|seed|p\.|pp\.|BY-NC-SA|Essays|Qwen|Llama|GPT-|DeBERTa-v|"
    r"Python|CUDA|torch|3\.1)\s*$", re.I)

# Thousands separators have to be part of the token, or "1,280 essays" is read as the number 280.
NUM = re.compile(r"(?<![\w.,])(\d{1,3}(?:,\d{3})+(?:\.\d+)?|\d+(?:\.\d+)?)(?![\w,]|\.\d)")
YEAR = re.compile(r"^(19|20)\d\d$")

# Corpus and design constants. They are facts about the study rather than results, so no results
# file holds them, and listing them here is cheaper than reading past them on every run.
CONSTANTS = {
    "640": "human essays sampled from BAWE",
    "627": "students in the source corpus",
    "901": "questions across the four arms, 270 + 204 + 243 + 184",
    "2000": "the target length asked of the generator, not a measurement",
    "2,000": "the target length asked of the generator, not a measurement",
    "1280": "essays in the detection corpus, 640 pairs",
    "1,280": "essays in the detection corpus, 640 pairs",
    "80": "essays per stratification cell",
    "402": "essays in the Persuasive Essays corpus",
    "200": "essays in the held-out detection test set",
    "50": "human essays per language group in the test set",
    "42": "the random seed",
    "28": "submission date in August",
    "8": "GB of VRAM",
}


def norm(val: str) -> str:
    """Normalise a written number for comparison. Prose writes 0.990 and 1.00 where a file holds
    0.99 and 1.0, and the difference is presentation rather than substance."""
    bare = val.replace(",", "")
    try:
        return f"{float(bare):g}"
    except ValueError:
        return bare


def skip(val, line, at) -> bool:
    """True for numbers that are not empirical claims about this project's results."""
    bare = val.replace(",", "")
    if bare in IGNORE_EXACT or norm(val) in IGNORE_EXACT or val in CONSTANTS or bare in CONSTANTS:
        return True
    if NEAR_STRUCTURAL.search(line[:at]):
        return True
    if YEAR.match(bare):
        # A year in a citation is followed by a close bracket or a semicolon, or preceded by "et
        # al.," or a surname and a comma.
        after = line[at + len(val): at + len(val) + 2]
        if after[:1] in (")", ";", ",", ":") or re.search(r"(et al\.|[A-Z][a-z]+)[,\s]+$", line[:at]):
            return True
    return False


def json_numbers():
    """Every numeric value anywhere in outputs/*.json, as a set of normalised strings.

    Values are recorded at several precisions because prose rounds: 0.9901 in the file is quoted as
    0.99 in the text, and 0.2007 as 20 percent."""
    seen: dict[str, set[str]] = {}

    def add(v, src):
        if not isinstance(v, (int, float)) or isinstance(v, bool):
            return
        # The number pattern does not capture a leading minus, so negatives have to be findable by
        # their magnitude or every correlation in the write-up reads as unsourced.
        forms = {f"{v:g}", f"{abs(v):g}"}
        if isinstance(v, float):
            for places in (0, 1, 2, 3, 4):
                forms.add(f"{round(v, places):g}")
                forms.add(f"{round(abs(v), places):g}")
            # Proportions are quoted as percentages at least as often as decimals.
            if 0 <= abs(v) <= 1:
                for places in (0, 1, 2):
                    forms.add(f"{round(abs(v) * 100, places):g}")
        else:
            forms.add(str(v))
        for f in forms:
            seen.setdefault(f, set()).add(src)

    def walk(o, src):
        if isinstance(o, dict):
            for v in o.values():
                walk(v, src)
        elif isinstance(o, list):
            for v in o:
                walk(v, src)
        else:
            add(o, src)

    for f in sorted(OUTPUTS.rglob("*.json")):
        try:
            walk(json.loads(f.read_text(encoding="utf8")), f.name)
        except Exception as e:
            print(f"  could not read {f.name}: {e}")
    return seen


def claims_from_chapters(include_lit=False):
    """(source, line number, number, context) for every number in the prose."""
    out = []
    for f in sorted(CHAPTERS.glob("*.md")):
        # Chapter 2 quotes other people's results by design and chapter 11 is the bibliography,
        # where every arXiv identifier and page range would be flagged. Neither belongs here.
        if f.name.startswith("11_"):
            continue
        if not include_lit and f.name.startswith("02_"):
            continue
        in_code = False
        for ln, line in enumerate(f.read_text(encoding="utf8").split("\n"), 1):
            if line.strip().startswith("```"):
                in_code = not in_code
                continue
            if in_code or SKIP_LINE.match(line):
                continue
            for m in NUM.finditer(line):
                val = m.group(1)
                if skip(val, line, m.start()):
                    continue
                ctx = line[max(0, m.start() - 46): m.end() + 34].strip()
                out.append((f.name, ln, val, ctx))
    return out


def claims_from_deck():
    from pptx import Presentation
    out = []
    prs = Presentation(DECK)
    for i, s in enumerate(prs.slides, 1):
        parts = [sh.text_frame.text for sh in s.shapes if sh.has_text_frame]
        parts.append(s.notes_slide.notes_text_frame.text)
        for text in parts:
            for line in text.split("\n"):
                for m in NUM.finditer(line):
                    val = m.group(1)
                    if skip(val, line, m.start()):
                        continue
                    ctx = line[max(0, m.start() - 46): m.end() + 34].strip()
                    out.append((f"slide {i}", 0, val, ctx))
    return out


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--deck", action="store_true", help="also check the presentation")
    ap.add_argument("--all", action="store_true", help="list matched numbers too")
    ap.add_argument("--lit", action="store_true",
                    help="include the literature review, whose numbers come from other papers")
    args = ap.parse_args()

    known = json_numbers()
    print(f"{len(known)} distinct numeric forms across {len(list(OUTPUTS.rglob('*.json')))} "
          f"results files\n")

    claims = claims_from_chapters(args.lit)
    if args.deck:
        claims += claims_from_deck()

    matched, unmatched = [], []
    for src, ln, val, ctx in claims:
        (matched if norm(val) in known else unmatched).append((src, ln, val, ctx))

    print(f"{len(claims)} numbers found in the text, {len(matched)} traceable to a results file, "
          f"{len(unmatched)} not.\n")

    by_source: dict[str, list] = {}
    for src, ln, val, ctx in unmatched:
        by_source.setdefault(src, []).append((ln, val, ctx))
    for src in sorted(by_source):
        print(f"--- {src}")
        for ln, val, ctx in by_source[src]:
            where = f":{ln}" if ln else ""
            print(f"  {val:>8}{where:>6}   {ctx}")
        print()

    if args.all:
        print("=" * 78)
        print("Matched, with the files that could be their source:")
        for src, ln, val, ctx in matched:
            files = sorted(known[val])
            tail = ", ".join(files[:3]) + (" ..." if len(files) > 3 else "")
            print(f"  {val:>8}  {src}:{ln}  <- {tail}")


if __name__ == "__main__":
    main()
